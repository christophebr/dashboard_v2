from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import io
import tempfile
import os
from data_processing.kpi_generation import (
    charge_entrant_sortant,
    charge_ticket,
    historique_scores_total,
    df_compute_ticket_appels_metrics,
    filtrer_par_periode,
    graph_activite,
    calculer_scores_equipe
)
import pandas as pd
from PIL import Image
import markdown
from bs4 import BeautifulSoup

# Couleurs pour les scores
SCORE_COLORS = {
    'low': RGBColor(255, 128, 128),     # Rouge clair
    'medium': RGBColor(255, 192, 0),     # Orange
    'high': RGBColor(146, 208, 80)       # Vert
}

# Couleurs pour les graphiques en ligne
LINE_COLORS = {
    'Archimede KESSI': '#1f77b4',        # Bleu
    'Morgane Vandenbussche': '#ff7f0e',  # Orange
    'Mourad HUMBLOT': '#2ca02c',         # Vert
    'Olivier Sainte-Rose': '#d62728'      # Rouge
}

def add_centered_picture(slide, image_path, width_inches=10):
    """Ajoute une image centrée dans la slide avec la largeur spécifiée"""
    # Dimensions standard d'une slide PowerPoint (13.33 x 7.5 pouces)
    SLIDE_WIDTH_INCHES = 13.33
    SLIDE_HEIGHT_INCHES = 7.5
    
    # Charger l'image pour obtenir ses proportions
    with Image.open(image_path) as img:
        img_width, img_height = img.size
        aspect_ratio = img_height / img_width
    
    # Calculer les dimensions finales
    picture_width = Inches(width_inches)
    picture_height = Inches(width_inches * aspect_ratio)
    
    # Calculer la position pour centrer
    left = Inches((SLIDE_WIDTH_INCHES - width_inches) / 2)
    top = Inches((SLIDE_HEIGHT_INCHES - (width_inches * aspect_ratio)) / 2)
    
    # Ajouter l'image
    slide.shapes.add_picture(image_path, left, top, width=picture_width)

def format_title(title_shape):
    """Formate le titre avec une police de taille 30"""
    title_shape.text_frame.paragraphs[0].font.size = Pt(30)

def seconds_to_time(seconds):
    """Convertit des secondes en format hh:mm:ss"""
    hours = int(seconds // 3600)
    minutes = int((seconds % 3600) // 60)
    seconds = int(seconds % 60)
    return f"{hours:02d}:{minutes:02d}:{seconds:02d}"

def create_powerpoint(kpis, df_support, periode):
    # Créer une nouvelle présentation
    prs = Presentation()
    
    # Slide 1 - Titre
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Dashboard Support"
    format_title(title)
    subtitle.text = f"Période : {periode}"
    
    # Slide 2 - KPIs principaux
    slide_layout = prs.slide_layouts[5]  # Layout avec titre uniquement
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "KPIs Principaux"
    format_title(title)
    
    # Convertir le temps moyen en format hh:mm:ss
    temps_moyen = seconds_to_time(float(kpis['temps_moy_appel']))
    
    # Positions des KPIs
    kpi_positions = [
        (1, 2, "Taux de service", f"{kpis['Taux_de_service']}%"),
        (5, 2, "Appels entrants par jour", str(kpis['Entrant'])),
        (1, 4, "Numéros uniques par jour", str(kpis['Numero_unique'])),
        (5, 4, "Temps moyen par appel", temps_moyen)
    ]
    
    # Ajouter les KPIs dans des rectangles
    for left, top, label, value in kpi_positions:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(3.5),
            Inches(1.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
        
        text_frame = shape.text_frame
        p1 = text_frame.paragraphs[0]
        p1.text = f"{label}\n{value}"
        p1.font.size = Pt(14)
        p1.font.color.rgb = RGBColor(0, 0, 0)  # Noir
        p1.alignment = PP_ALIGN.CENTER

    # Configuration commune pour les graphiques
    chart_layout = {
        'template': 'plotly_white',
        'showlegend': True,
        'legend': {
            'orientation': 'h',
            'yanchor': 'bottom',
            'y': 1.02,
            'xanchor': 'right',
            'x': 1
        },
        'margin': {'l': 40, 'r': 40, 't': 60, 'b': 40},
        'xaxis': {
            'showgrid': True,
            'gridwidth': 1,
            'gridcolor': 'lightgray'
        },
        'yaxis': {
            'showgrid': True,
            'gridwidth': 1,
            'gridcolor': 'lightgray',
            'rangemode': 'tozero'
        },
        'width': 800,
        'height': 450,
        'font': {'size': 10}
    }
    
    # Slide 3 - Graphique d'activité
    img_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(img_slide_layout)
    title = slide.shapes.title
    title.text = "Activité journalière"
    format_title(title)
    
    # Générer et ajouter le graphique d'activité
    fig_activite = graph_activite(df_support)
    fig_activite.update_layout(chart_layout)
    img_bytes = io.BytesIO()
    fig_activite.write_image(img_bytes, format="png")
    img_bytes.seek(0)
    left = (prs.slide_width - Inches(8)) / 2
    slide.shapes.add_picture(img_bytes, left, Inches(1.5), width=Inches(8))
    
    # Slides suivants - Autres graphiques
    for fig in [kpis['charge_affid_stellair_%'], kpis['charge_affid_stellair_v'], kpis['evo_appels_tickets']]:
        img_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(img_slide_layout)
        title = slide.shapes.title
        format_title(title)
        
        # Mettre à jour et ajouter le graphique
        fig.update_layout(chart_layout)
        img_bytes = io.BytesIO()
        fig.write_image(img_bytes, format="png")
        img_bytes.seek(0)
        left = (prs.slide_width - Inches(8)) / 2
        slide.shapes.add_picture(img_bytes, left, Inches(1.5), width=Inches(8))
    
    # Sauvegarder la présentation en mémoire
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    return pptx_io

def create_powerpoint_agents(df_support, df_tickets, periode_selectbox, df_evaluation_filtre=None):
    print("Démarrage de la génération de la présentation")
    
    # Créer une nouvelle présentation
    prs = Presentation()
    
    try:
        print("Filtrage des données par période")
        # Filtrer les données selon la période sélectionnée
        df_support = filtrer_par_periode(df_support.copy(), periode_selectbox)
        df_tickets = filtrer_par_periode(df_tickets.copy(), periode_selectbox)

        print("Création du slide de titre")
        # Slide 1 - Titre
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Dashboard Agents N1"
        format_title(title)  # Formater le titre
        subtitle.text = f"Période : {periode_selectbox}"

        print("Calcul des scores des agents")
        # Générer le tableau des scores
        agents_n1 = ['Olivier Sainte-Rose', 'Mourad HUMBLOT', 'Archimede KESSI', 'Morgane Vandenbussche']
        agents_n1_tickets = agents_n1 + ['Frederic SAUVAN']

        print("Calcul des métriques tickets/appels")
        # Utiliser df_compute_ticket_appels_metrics pour la cohérence avec le dashboard
        df_conforme, _ = df_compute_ticket_appels_metrics(agents_n1, df_tickets, df_support)

        # Garder seulement les colonnes nécessaires dans l'ordre souhaité
        df_conforme = df_conforme[['Agent', 'score_performance', "Nombre d'appel traité", 
                                 'Nombre de ticket traité', '% tickets', '% appels',
                                 '% appel entrant agent']]

        # Slide 2 - Tableau des scores
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Scores des agents"
        format_title(title_shape)  # Formater le titre
        
        # Ajouter une description des règles de score
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(0.5))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = "Règles de score : > 70% (vert), 60-70% (orange), < 60% (rouge)"
        p.font.size = Pt(10)
        p.font.italic = True
        
        # Créer le tableau
        rows = len(df_conforme) + 1
        cols = len(df_conforme.columns)
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(11)
        height = Inches(0.2 * rows)
        
        table = shapes.add_table(rows, cols, left, top, width, height).table
        
        # Définir les largeurs de colonnes
        col_widths = [2.0, 1.5, 1.5, 1.5, 1.5, 1.5, 1.5]  # Total = 11 pouces
        for i, width in enumerate(col_widths):
            if i < len(table.columns):
                table.columns[i].width = Inches(width)
        
        # En-têtes
        header_color = RGBColor(217, 217, 217)  # Gris clair
        for i, column in enumerate(df_conforme.columns):
            cell = table.cell(0, i)
            cell.text = str(column)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(9)
            paragraph.alignment = PP_ALIGN.CENTER
        
        # Données avec alternance de couleurs
        row_color = RGBColor(242, 242, 242)  # Gris très clair
        for i, row in enumerate(df_conforme.values):
            for j, value in enumerate(row):
                cell = table.cell(i + 1, j)
                # Appliquer la couleur de fond alternée
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = row_color
                
                if isinstance(value, float):
                    cell.text = f"{value:.2f}"
                else:
                    cell.text = str(value)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.alignment = PP_ALIGN.CENTER
                
                # Colorer les scores selon les seuils
                if j == 1:  # Colonne du score_performance
                    if isinstance(value, (int, float)):
                        if value >= 70:
                            paragraph.font.color.rgb = SCORE_COLORS['high']
                        elif value >= 60:
                            paragraph.font.color.rgb = SCORE_COLORS['medium']
                        else:
                            paragraph.font.color.rgb = SCORE_COLORS['low']

        # Configuration commune pour les graphiques en ligne
        line_chart_layout = {
            'template': 'plotly_white',
            'showlegend': True,
            'legend': {
                'orientation': 'h',
                'yanchor': 'bottom',
                'y': 1.02,
                'xanchor': 'right',
                'x': 1
            },
            'margin': {'l': 40, 'r': 40, 't': 60, 'b': 40},  # Marges réduites
            'xaxis': {
                'showgrid': True,
                'gridwidth': 1,
                'gridcolor': 'lightgray'
            },
            'yaxis': {
                'showgrid': True,
                'gridwidth': 1,
                'gridcolor': 'lightgray',
                'rangemode': 'tozero'
            },
            'width': 800,  # Largeur réduite
            'height': 450,  # Hauteur réduite
            'font': {'size': 10}  # Taille de police réduite
        }

        # Configuration commune pour les camemberts
        pie_chart_layout = {
            'template': 'plotly_white',
            'showlegend': True,
            'legend': {
                'orientation': 'h',
                'yanchor': 'bottom',
                'y': 1.02,
                'xanchor': 'right',
                'x': 1
            },
            'margin': {'l': 20, 'r': 20, 't': 60, 'b': 20},  # Marges très réduites
            'width': 600,
            'height': 500,
            'font': {'size': 10}
        }

        print("Création des graphiques d'activité téléphonique")
        # Graphiques d'activité téléphonique
        fig_line_entrant, fig_pie_entrant = charge_entrant_sortant(df_support, agents_n1)
        
        # Appliquer les couleurs aux lignes pour les graphiques d'appels
        if hasattr(fig_line_entrant.data[0], 'line'):
            for trace in fig_line_entrant.data:
                if trace.name in LINE_COLORS:
                    trace.line.color = LINE_COLORS[trace.name]

        # Graphiques d'activité tickets
        fig_line_ticket, fig_pie_ticket = charge_ticket(df_tickets, agents_n1_tickets)
        
        # Appliquer les couleurs aux lignes pour les graphiques de tickets
        if hasattr(fig_line_ticket.data[0], 'line'):
            for trace in fig_line_ticket.data:
                if trace.name in LINE_COLORS:
                    trace.line.color = LINE_COLORS[trace.name]

        # Générer et configurer le graphique d'historique
        fig_historique = historique_scores_total(agents_n1, df_tickets, df_support)
        
        # Appliquer les couleurs aux lignes pour l'historique
        if hasattr(fig_historique.data[0], 'line'):
            for trace in fig_historique.data:
                if trace.name in LINE_COLORS:
                    trace.line.color = LINE_COLORS[trace.name]

        # Appliquer la configuration aux graphiques en ligne
        fig_line_entrant.update_layout(line_chart_layout)
        fig_line_ticket.update_layout(line_chart_layout)
        fig_historique.update_layout(line_chart_layout)

        # Appliquer la configuration aux camemberts
        fig_pie_entrant.update_layout(pie_chart_layout)
        fig_pie_ticket.update_layout(pie_chart_layout)

        # Slide 3 - Graphique en ligne des appels
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Volume d'appels entrants par agent"
        format_title(title)  # Formater le titre
        
        # Sauvegarder et ajouter le graphique en ligne
        img_bytes = io.BytesIO()
        fig_line_entrant.write_image(img_bytes, format="png")
        img_bytes.seek(0)
        left = (prs.slide_width - Inches(8)) / 2  # Centrage automatique
        slide.shapes.add_picture(img_bytes, left, Inches(1.5), width=Inches(8))

        # Slide 4 - Graphique camembert des appels
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Répartition des appels entrants par agent"
        format_title(title)  # Formater le titre
        
        # Sauvegarder et ajouter le graphique camembert
        img_bytes = io.BytesIO()
        fig_pie_entrant.write_image(img_bytes, format="png")
        img_bytes.seek(0)
        left = (prs.slide_width - Inches(6)) / 2  # Centrage automatique
        slide.shapes.add_picture(img_bytes, left, Inches(1.5), width=Inches(6))

        # Slide 5 - Graphique en ligne des tickets
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Volume de tickets par agent"
        format_title(title)  # Formater le titre
        
        # Sauvegarder et ajouter le graphique en ligne
        img_bytes = io.BytesIO()
        fig_line_ticket.write_image(img_bytes, format="png")
        img_bytes.seek(0)
        left = (prs.slide_width - Inches(8)) / 2  # Centrage automatique
        slide.shapes.add_picture(img_bytes, left, Inches(1.5), width=Inches(8))
        
        # Slide 6 - Graphique camembert des tickets
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Répartition des tickets par agent"
        format_title(title)  # Formater le titre
        
        # Sauvegarder et ajouter le graphique camembert
        img_bytes = io.BytesIO()
        fig_pie_ticket.write_image(img_bytes, format="png")
        img_bytes.seek(0)
        left = (prs.slide_width - Inches(6)) / 2  # Centrage automatique
        slide.shapes.add_picture(img_bytes, left, Inches(1.5), width=Inches(6))

        # Slide 7 - Historique des scores
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Historique des scores globaux"
        format_title(title)  # Formater le titre
        
        # Sauvegarder et ajouter le graphique
        img_bytes = io.BytesIO()
        fig_historique.write_image(img_bytes, format="png")
        img_bytes.seek(0)
        left = (prs.slide_width - Inches(8)) / 2  # Centrage automatique
        slide.shapes.add_picture(img_bytes, left, Inches(1.5), width=Inches(8))

        # Slide 8 - Évaluations managériales (si disponible)
        if df_evaluation_filtre is not None and not df_evaluation_filtre.empty:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = "Évaluations managériales"
            format_title(title)

            # Calculer la moyenne des scores par agent et par trimestre
            df_summary = df_evaluation_filtre.groupby(['agent', 'quarter'])['average_score'].mean().reset_index()
            moyenne_score = round(df_summary['average_score'].mean(), 2)

            # Créer le tableau
            rows = len(df_summary) + 1  # +1 pour l'en-tête
            cols = len(df_summary.columns)
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(11)
            height = Inches(0.3 * rows)

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            # En-têtes
            header_color = RGBColor(217, 217, 217)  # Gris clair
            for i, column in enumerate(df_summary.columns):
                cell = table.cell(0, i)
                cell.text = str(column)
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.bold = True
                paragraph.font.size = Pt(10)
                paragraph.alignment = PP_ALIGN.CENTER

            # Données
            for i, row in enumerate(df_summary.values):
                for j, value in enumerate(row):
                    cell = table.cell(i + 1, j)
                    if isinstance(value, float):
                        cell.text = f"{value:.2f}"
                    else:
                        cell.text = str(value)
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(9)
                    paragraph.alignment = PP_ALIGN.CENTER

            # Ajouter le score moyen
            score_text = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(0.5))
            score_text.text_frame.text = f"Score moyen des évaluations : {moyenne_score}"
            score_text.text_frame.paragraphs[0].font.bold = True

        # Slide 4 - Evolution temps de réponse aux tickets
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Évolution du temps de réponse aux tickets (6 mois)"
        format_title(title)
        img_bytes = io.BytesIO()
        # Mettre la légende en haut, horizontale et centrée
        graph_temps_reponse_6m.update_layout(
            legend=dict(orientation="h", yanchor="bottom", y=1.08, xanchor="center", x=0.5)
        )
        # Exporter l'image avec la bonne taille
        graph_temps_reponse_6m.write_image(img_bytes, format="png", width=9.06*96, height=5.51*96)  # 96 dpi
        img_bytes.seek(0)
        width = Inches(9.06)
        height = Inches(5.51)
        left = (prs.slide_width - width) / 2
        slide.shapes.add_picture(img_bytes, left, Inches(1.2), width=width, height=height)

        # Sauvegarder la présentation en mémoire
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        
        return pptx_io
        
    except Exception as e:
        print(f"Erreur lors de la génération de la présentation : {str(e)}")
        raise 

def create_powerpoint_stellair_report(df_support_6m, df_tickets_6m, kpis_6m, df_support_3m, df_tickets_3m, kpis_3m, graph_activite_6m, graph_temps_reponse_6m, evo_appels_tickets_6m):
    """
    Génère un rapport PowerPoint Stellair avec :
    - Slide 1 : KPI 6 mois
    - Slide 2 : KPI 3 mois
    - Slide 3 : Graphique activité 6 mois
    - Slide 4 : Evolution temps de réponse aux tickets 6 mois
    - Slide 5 : Evolution hebdo appels entrants + tickets 6 mois
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    import io
    
    prs = Presentation()
    
    # Slide 1 - KPIs 6 mois
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Support Stellair - 6 mois : KPIs principaux"
    format_title(title)
    kpi_labels = [
        ("Taux de service", f"{kpis_6m['Taux_de_service']}%"),
        ("Appels entrant / Jour", str(kpis_6m['Entrant'])),
        ("Numéros uniques / Jour", str(kpis_6m['Numero_unique'])),
        ("Entrants vs Tickets", f"{round(kpis_6m['activite_appels_pourcentage']*100,2)}% / {round(kpis_6m['activite_tickets_pourcentage']*100,2)}%"),
        ("Temps de réponse moyen aux tickets (h:min)", f"{int(kpis_6m['moyenne_temps_reponse'])}:{int((kpis_6m['moyenne_temps_reponse']%1)*60):02d}")
    ]
    for i, (label, value) in enumerate(kpi_labels):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5 + (i%3)*3.2),  # plus rapproché
            Inches(1.2 + (i//3)*1.6), # plus haut sur la slide
            Inches(3),                # moins large
            Inches(1.1)               # moins haut
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
        text_frame = shape.text_frame
        p1 = text_frame.paragraphs[0]
        p1.text = f"{label}\n{value}"
        p1.font.size = Pt(16)
        p1.font.color.rgb = RGBColor(0, 0, 0)
        p1.alignment = PP_ALIGN.CENTER
    
    # Slide 2 - KPIs 3 mois
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Support Stellair - 3 mois : KPIs principaux"
    format_title(title)
    kpi_labels = [
        ("Taux de service", f"{kpis_3m['Taux_de_service']}%"),
        ("Appels entrant / Jour", str(kpis_3m['Entrant'])),
        ("Numéros uniques / Jour", str(kpis_3m['Numero_unique'])),
        ("Entrants vs Tickets", f"{round(kpis_3m['activite_appels_pourcentage']*100,2)}% / {round(kpis_3m['activite_tickets_pourcentage']*100,2)}%"),
        ("Temps de réponse moyen aux tickets (h:min)", f"{int(kpis_3m['moyenne_temps_reponse'])}:{int((kpis_3m['moyenne_temps_reponse']%1)*60):02d}")
    ]
    for i, (label, value) in enumerate(kpi_labels):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5 + (i%3)*3.2),
            Inches(1.2 + (i//3)*1.6),
            Inches(3),
            Inches(1.1)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
        text_frame = shape.text_frame
        p1 = text_frame.paragraphs[0]
        p1.text = f"{label}\n{value}"
        p1.font.size = Pt(16)
        p1.font.color.rgb = RGBColor(0, 0, 0)
        p1.alignment = PP_ALIGN.CENTER
    
    # Slide 3 - Graphique activité 6 mois
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Support Stellair - 6 mois : Activité"
    format_title(title)
    img_bytes = io.BytesIO()
    graph_activite_6m.write_image(img_bytes, format="png")
    img_bytes.seek(0)
    left = (prs.slide_width - Inches(8)) / 2
    slide.shapes.add_picture(img_bytes, left, Inches(1.5), width=Inches(8))
    
    # Slide 4 - Evolution temps de réponse aux tickets
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Évolution du temps de réponse aux tickets (6 mois)"
    format_title(title)
    img_bytes = io.BytesIO()
    # Mettre la légende en haut, horizontale et centrée
    graph_temps_reponse_6m.update_layout(
        legend=dict(orientation="h", yanchor="bottom", y=1.08, xanchor="center", x=0.5)
    )
    # Exporter l'image avec la bonne taille
    graph_temps_reponse_6m.write_image(img_bytes, format="png", width=9.06*96, height=5.51*96)  # 96 dpi
    img_bytes.seek(0)
    width = Inches(9.06)
    height = Inches(5.51)
    left = (prs.slide_width - width) / 2
    slide.shapes.add_picture(img_bytes, left, Inches(1.2), width=width, height=height)
    
    # Slide 5 - Evolution hebdo appels entrants + tickets
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Évolution hebdomadaire : Appels entrants + Tickets (6 mois)"
    format_title(title)
    img_bytes = io.BytesIO()
    evo_appels_tickets_6m.write_image(img_bytes, format="png")
    img_bytes.seek(0)
    left = (prs.slide_width - Inches(8)) / 2
    slide.shapes.add_picture(img_bytes, left, Inches(1.5), width=Inches(8))
    
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io 

def create_powerpoint_agents_report(df_scores_6m, df_scores_3m, fig_tel_6m, fig_ticket_6m, markdown_methodo):
    """
    Génère un rapport PowerPoint Agents avec :
    - Slide 1 : Titre, texte méthodologie scoring, tableau scoring 6 mois, tableau scoring 3 mois
    - Slide 2 : Graphique activité téléphonique (6 mois)
    - Slide 3 : Graphique activité tickets (6 mois)
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    import io
    import pandas as pd
    import markdown
    from bs4 import BeautifulSoup

    prs = Presentation()

    # Slide 1 - Méthodologie et tableaux
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Méthodologie scoring et résultats"
    format_title(title)

    # Texte méthodologie (markdown -> texte simple)
    html = markdown.markdown(markdown_methodo)
    soup = BeautifulSoup(html, "html.parser")
    methodo_text = soup.get_text(separator="\n")
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8.5), Inches(1.2))
    tf = textbox.text_frame
    tf.text = methodo_text
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Tableau scoring 6 mois
    rows, cols = df_scores_6m.shape
    table = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(2.6), Inches(8.5), Inches(1.2+0.3*rows)).table
    for j, col in enumerate(df_scores_6m.columns):
        table.cell(0, j).text = str(col)
    for i in range(rows):
        for j in range(cols):
            table.cell(i+1, j).text = str(df_scores_6m.iloc[i, j])
    for cell in table.iter_cells():
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    # Tableau scoring 3 mois (en dessous)
    rows2, cols2 = df_scores_3m.shape
    table2 = slide.shapes.add_table(rows2+1, cols2, Inches(0.5), Inches(2.7+0.3*rows), Inches(8.5), Inches(1.2+0.3*rows2)).table
    for j, col in enumerate(df_scores_3m.columns):
        table2.cell(0, j).text = str(col)
    for i in range(rows2):
        for j in range(cols2):
            table2.cell(i+1, j).text = str(df_scores_3m.iloc[i, j])
    for cell in table2.iter_cells():
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    # Slide 2 - Graphique activité téléphonique (6 mois)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Activité téléphonique (6 mois)"
    format_title(title)
    img_bytes = io.BytesIO()
    fig_tel_6m.write_image(img_bytes, format="png", width=9.06*96, height=5.51*96)
    img_bytes.seek(0)
    width = Inches(9.06)
    height = Inches(5.51)
    left = (prs.slide_width - width) / 2
    slide.shapes.add_picture(img_bytes, left, Inches(1.2), width=width, height=height)

    # Slide 3 - Graphique activité tickets (6 mois)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Activité tickets (6 mois)"
    format_title(title)
    img_bytes = io.BytesIO()
    fig_ticket_6m.write_image(img_bytes, format="png", width=9.06*96, height=5.51*96)
    img_bytes.seek(0)
    width = Inches(9.06)
    height = Inches(5.51)
    left = (prs.slide_width - width) / 2
    slide.shapes.add_picture(img_bytes, left, Inches(1.2), width=width, height=height)

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io 